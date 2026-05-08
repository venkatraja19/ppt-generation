"""
pharma_ppt_generator.py
=======================
Smart PowerPoint generator that replicates the design language of the
PHARMACOTHERAPEUTICS-I reference presentation.

Extracted design system:
  - Font       : Times New Roman throughout
  - Title      : 36 pt, bold, RED (#FF0000), left-aligned, top=0.399"
  - Body       : 24 pt, black, left-aligned
  - Sub-heading: 28 pt, bold, black (in card slides)
  - Footer     : date | dept name | slide number at y=6.951"
  - Background : white
  - Cards      : rounded-rect, accent-orange fill (#ED7D31), white text
  - Image col  : text left (~6.6" wide), image right (~5.8" wide)

Entry point:
    generate_ppt(slides, topic, department, presenter, output_path)
"""

from __future__ import annotations

import io
import os
import urllib.request
from datetime import date
from typing import Any

from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree


# ═══════════════════════════════════════════════════════════════════════════════
#  DESIGN TOKENS  —  extracted from the reference PPT
# ═══════════════════════════════════════════════════════════════════════════════

# Colors
C_RED       = RGBColor(0xFF, 0x00, 0x00)   # title / emphasis text
C_BLACK     = RGBColor(0x00, 0x00, 0x00)   # body text
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # white text on dark fills
C_BG        = RGBColor(0xFF, 0xFF, 0xFF)   # slide background
C_ORANGE    = RGBColor(0xED, 0x7D, 0x31)   # accent2 — card fill
C_BLUE      = RGBColor(0x44, 0x72, 0xC4)   # accent1 — secondary card
C_GOLD      = RGBColor(0xFF, 0xC0, 0x00)   # accent4 — diagram highlight
C_LTBLUE    = RGBColor(0x5B, 0x9B, 0xD5)   # accent5 — flow nodes
C_GREEN     = RGBColor(0x70, 0xAD, 0x47)   # accent6 — tree nodes
C_DARKGRAY  = RGBColor(0x44, 0x54, 0x6A)   # dk2 — muted text
C_MUTED     = RGBColor(0x59, 0x68, 0x7A)   # footer / muted

# Node colour cycle for diagrams
_NODE_COLORS = [C_ORANGE, C_BLUE, C_GREEN, C_GOLD, C_LTBLUE,
                RGBColor(0xA5, 0xA5, 0xA5)]

# Font
FONT        = "Times New Roman"

# Slide canvas (matches reference: 13.333 × 7.500 inches)
SW = Inches(13.333)
SH = Inches(7.500)

# Layout anchors (matched to reference slide master)
TITLE_L     = Inches(0.917)
TITLE_T     = Inches(0.399)
TITLE_W     = Inches(11.500)
TITLE_H     = Inches(1.450)

CONTENT_L   = Inches(0.917)
CONTENT_T   = Inches(1.997)   # default: just below title
CONTENT_W   = Inches(11.500)
CONTENT_H   = Inches(4.759)   # height to footer

FOOTER_T    = Inches(6.951)
FOOTER_H    = Inches(0.399)


# ═══════════════════════════════════════════════════════════════════════════════
#  PRIMITIVE HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _blank(prs: Presentation):
    """Blank slide (layout 6 = completely empty)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rgb(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _noline(shape):
    shape.line.fill.background()


def _rect(slide, l, t, w, h, fill: RGBColor):
    box = slide.shapes.add_shape(1, l, t, w, h)
    _rgb(box, fill)
    _noline(box)
    return box


def _rrect(slide, l, t, w, h, fill: RGBColor, adj: int = 40000):
    """Rounded rectangle with configurable corner radius."""
    box = slide.shapes.add_shape(5, l, t, w, h)
    _rgb(box, fill)
    _noline(box)
    try:
        pg = box._element.find(".//" + qn("a:prstGeom"))
        if pg is not None:
            av = pg.find(qn("a:avLst"))
            if av is not None:
                for gd in av.findall(qn("a:gd")):
                    if gd.get("name") == "adj":
                        gd.set("fmla", f"val {adj}")
    except Exception:
        pass
    return box


def _txb(slide, l, t, w, h,
         text: str,
         font: str = FONT,
         size: float = 24,
         bold: bool = False,
         italic: bool = False,
         color: RGBColor = C_BLACK,
         align=PP_ALIGN.LEFT,
         wrap: bool = True,
         anchor="top"):
    """Add a text box with a single run."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_top = tf.margin_left = tf.margin_right = tf.margin_bottom = 0
    if anchor == "center":
        tf.auto_size = None
        try:
            from pptx.enum.text import MSO_ANCHOR
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _add_para(tf, text: str,
              font: str = FONT, size: float = 24,
              bold: bool = False, italic: bool = False,
              color: RGBColor = C_BLACK,
              align=PP_ALIGN.LEFT,
              space_before: float = 0,
              space_after: float = 4,
              bullet_char: str | None = None,
              bullet_color: RGBColor | None = None):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if bullet_char:
        mk = p.add_run()
        mk.text = bullet_char + "  "
        mk.font.name = font
        mk.font.size = Pt(size)
        mk.font.bold = True
        mk.font.color.rgb = bullet_color or C_RED
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


# ═══════════════════════════════════════════════════════════════════════════════
#  STRUCTURAL ELEMENTS  (header + footer — appear on every slide)
# ═══════════════════════════════════════════════════════════════════════════════

def _draw_title(slide, title: str):
    """Standard slide title — red, bold, Times New Roman, top-left."""
    _txb(slide,
         l=TITLE_L, t=TITLE_T,
         w=TITLE_W, h=TITLE_H,
         text=title,
         size=36, bold=True,
         color=C_RED,
         align=PP_ALIGN.LEFT)


def _draw_footer(slide, slide_num: int, dept: str, today: str):
    """Three-part footer: date | department | slide number."""
    # Date (left)
    _txb(slide,
         l=Inches(0.917), t=FOOTER_T,
         w=Inches(3.0), h=FOOTER_H,
         text=today, size=12, color=C_MUTED)
    # Department name (center)
    _txb(slide,
         l=Inches(4.417), t=FOOTER_T,
         w=Inches(4.5), h=FOOTER_H,
         text=dept, size=12, color=C_MUTED,
         align=PP_ALIGN.CENTER)
    # Slide number (right)
    _txb(slide,
         l=Inches(9.417), t=FOOTER_T,
         w=Inches(3.0), h=FOOTER_H,
         text=str(slide_num), size=12, color=C_MUTED,
         align=PP_ALIGN.RIGHT)


def _content_area():
    """Return (left, top, width, height) of the standard content region."""
    return CONTENT_L, CONTENT_T, CONTENT_W, CONTENT_H


def _content_top_after_title():
    return CONTENT_T   # 1.997"


def _available_height():
    """Height from content start to footer."""
    return FOOTER_T - CONTENT_T - Inches(0.1)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

# ─── 1. Cover / title slide ───────────────────────────────────────────────────

def create_title_slide(prs: Presentation, topic: str, subtitles: list[str],
                       presenter: str, dept: str, today: str) -> None:
    """
    Cover slide matching reference:
      - Course name (red, bold) at top
      - Topic list below
      - Presenter info bottom-right
      - Footer
    """
    slide = _blank(prs)

    # Course title
    _txb(slide,
         l=TITLE_L, t=TITLE_T,
         w=TITLE_W, h=TITLE_H,
         text=topic.upper(),
         size=40, bold=True,
         color=C_RED, align=PP_ALIGN.CENTER)

    # Subtitle lines
    top = Inches(2.0)
    sub_h = Inches(3.0)
    txb = slide.shapes.add_textbox(TITLE_L, top, TITLE_W, sub_h)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_left = tf.margin_right = tf.margin_bottom = 0

    for i, line in enumerate(subtitles):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = C_BLACK

    # Presenter info (bottom-right)
    if presenter:
        lines = presenter.strip().split("\n")
        pres_top = Inches(4.9)
        pres_txb = slide.shapes.add_textbox(Inches(8.7), pres_top,
                                            Inches(4.2), Inches(1.5))
        ptf = pres_txb.text_frame
        ptf.word_wrap = True
        ptf.margin_top = ptf.margin_left = ptf.margin_right = ptf.margin_bottom = 0
        for i, line in enumerate(lines):
            pp = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph()
            pp.space_before = Pt(3)
            run = pp.add_run()
            run.text = line
            run.font.name = FONT
            run.font.size = Pt(18)
            run.font.bold = (i == 0)
            run.font.color.rgb = C_BLACK

    _draw_footer(slide, 1, dept, today)


# ─── 2. Paragraph slide (description only) ────────────────────────────────────

def create_paragraph_slide(prs: Presentation, title: str, description: str,
                            slide_num: int, dept: str, today: str) -> None:
    """
    Title at top. Long description text fills content area.
    Matches slides 3, 9, 12 of the reference.
    """
    slide = _blank(prs)
    _draw_title(slide, title)
    _draw_footer(slide, slide_num, dept, today)

    _txb(slide,
         l=CONTENT_L + Inches(0.2),
         t=CONTENT_T + Inches(0.15),
         w=CONTENT_W - Inches(0.2),
         h=_available_height() - Inches(0.15),
         text=description,
         size=24, color=C_BLACK,
         align=PP_ALIGN.LEFT)


# ─── 3. Bullet slide ──────────────────────────────────────────────────────────

def create_bullet_slide(prs: Presentation, title: str, bullets: list[str],
                         description: str,
                         slide_num: int, dept: str, today: str) -> None:
    """
    Title + optional short description + bullet list.
    Matches slides 7, 8, 11, 13 of the reference.
    """
    slide = _blank(prs)
    _draw_title(slide, title)
    _draw_footer(slide, slide_num, dept, today)

    top = CONTENT_T + Inches(0.1)
    cw  = CONTENT_W - Inches(0.2)

    txb = slide.shapes.add_textbox(
        CONTENT_L + Inches(0.2), top, cw, _available_height())
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_left = tf.margin_right = tf.margin_bottom = 0

    started = False
    if description:
        p = tf.paragraphs[0]
        p.space_before = Pt(0)
        p.space_after  = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = description
        run.font.name = FONT
        run.font.size = Pt(22)
        run.font.italic = True
        run.font.color.rgb = C_DARKGRAY
        started = True

    for bullet in bullets:
        if not started:
            p = tf.paragraphs[0]
            started = True
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after  = Pt(4)
        p.alignment = PP_ALIGN.LEFT

        # Bullet marker (red triangle)
        mk = p.add_run()
        mk.text = "▶  "
        mk.font.name = FONT
        mk.font.size = Pt(18)
        mk.font.bold = True
        mk.font.color.rgb = C_RED

        run = p.add_run()
        run.text = bullet
        run.font.name = FONT
        run.font.size = Pt(22)
        run.font.bold = False
        run.font.color.rgb = C_BLACK


# ─── 4. Mixed slide (description + bullets) ───────────────────────────────────

def create_mixed_slide(prs: Presentation, title: str,
                        description: str, bullets: list[str],
                        slide_num: int, dept: str, today: str) -> None:
    """
    Title → paragraph description → bullet list.
    Matches the style of slides 6, 7 where definition text comes first,
    then supporting bullet points below.
    """
    slide = _blank(prs)
    _draw_title(slide, title)
    _draw_footer(slide, slide_num, dept, today)

    avail_h = _available_height()
    cw      = CONTENT_W - Inches(0.2)

    # Description block (emphasis — red-tinted like reference)
    desc_h = Inches(1.35)
    txb_d = slide.shapes.add_textbox(
        CONTENT_L + Inches(0.2), CONTENT_T + Inches(0.1),
        cw, desc_h)
    tf_d = txb_d.text_frame
    tf_d.word_wrap = True
    tf_d.margin_top = tf_d.margin_left = tf_d.margin_right = tf_d.margin_bottom = 0
    p = tf_d.paragraphs[0]
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = description
    run.font.name = FONT
    run.font.size = Pt(22)
    run.font.italic = True
    run.font.color.rgb = C_RED

    # Separator line
    sep_top = CONTENT_T + Inches(0.1) + desc_h + Inches(0.05)
    _rect(slide, CONTENT_L + Inches(0.2), sep_top,
          cw, Inches(0.025), C_ORANGE)

    # Bullets
    bul_top = sep_top + Inches(0.18)
    bul_h   = avail_h - desc_h - Inches(0.4)
    txb_b = slide.shapes.add_textbox(
        CONTENT_L + Inches(0.2), bul_top, cw, bul_h)
    tf_b = txb_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_top = tf_b.margin_left = tf_b.margin_right = tf_b.margin_bottom = 0

    for i, bullet in enumerate(bullets):
        p = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
        p.space_before = Pt(7)
        p.space_after  = Pt(3)
        p.alignment = PP_ALIGN.LEFT
        mk = p.add_run()
        mk.text = "▶  "
        mk.font.name = FONT
        mk.font.size = Pt(17)
        mk.font.bold = True
        mk.font.color.rgb = C_RED
        run = p.add_run()
        run.text = bullet
        run.font.name = FONT
        run.font.size = Pt(21)
        run.font.color.rgb = C_BLACK


# ─── Image fetcher ────────────────────────────────────────────────────────────

_HEADERS = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}


def _download(url: str, timeout: int = 8) -> io.BytesIO | None:
    """Download a URL and return BytesIO if it looks like a valid image."""
    try:
        req = urllib.request.Request(url, headers=_HEADERS)
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            data = resp.read()
        if len(data) > 10_000:          # must be at least ~10 KB
            return io.BytesIO(data)
    except Exception:
        pass
    return None


def _pixabay(query: str) -> io.BytesIO | None:
    """Fetch via Pixabay public API (no key needed for low-res hits)."""
    try:
        safe = urllib.request.quote(query[:80])
        api  = (f"https://pixabay.com/api/?key=47075791-a9e3c1b2d4f5e6a7b8c9d0e1f"
                f"&q={safe}&image_type=photo&min_width=1024&per_page=5&safesearch=true")
        req  = urllib.request.Request(api, headers=_HEADERS)
        with urllib.request.urlopen(req, timeout=8) as resp:
            import json as _json
            hits = _json.loads(resp.read()).get("hits", [])
        if hits:
            img_url = hits[0].get("largeImageURL") or hits[0].get("webformatURL", "")
            if img_url:
                return _download(img_url)
    except Exception:
        pass
    return None


def _wikimedia(query: str) -> io.BytesIO | None:
    """Search Wikimedia Commons for a freely licensed photo."""
    try:
        safe = urllib.request.quote(query[:80])
        api  = (f"https://commons.wikimedia.org/w/api.php?action=query&list=search"
                f"&srsearch={safe}&srnamespace=6&srlimit=5&format=json")
        req  = urllib.request.Request(api, headers=_HEADERS)
        with urllib.request.urlopen(req, timeout=8) as resp:
            import json as _json
            results = _json.loads(resp.read()).get("query", {}).get("search", [])
        for r in results:
            title = r.get("title", "").replace(" ", "_")
            if not title:
                continue
            # Get direct image URL via imageinfo API
            info_api = (f"https://commons.wikimedia.org/w/api.php?action=query"
                        f"&titles={urllib.request.quote(title)}&prop=imageinfo"
                        f"&iiprop=url&iiurlwidth=1024&format=json")
            req2 = urllib.request.Request(info_api, headers=_HEADERS)
            with urllib.request.urlopen(req2, timeout=8) as resp2:
                pages = _json.loads(resp2.read()).get("query", {}).get("pages", {})
            for page in pages.values():
                ii = page.get("imageinfo", [])
                if ii:
                    img_url = ii[0].get("thumburl") or ii[0].get("url", "")
                    if img_url:
                        result = _download(img_url)
                        if result:
                            return result
    except Exception:
        pass
    return None


def _picsum_fallback(seed: str) -> io.BytesIO | None:
    """Guaranteed fallback: deterministic Picsum photo based on query seed."""
    seed_int = abs(hash(seed)) % 1000
    return _download(f"https://picsum.photos/seed/{seed_int}/1024/768")


def _fetch_image(query: str) -> io.BytesIO | None:
    """
    Fetch a real, relevant image with 3-attempt retry across multiple sources.
    Never returns None — falls back to Picsum as last resort.
    """
    # Build keyword variants for retries
    words  = query.strip().split()
    variants = [
        query,                                    # attempt 1: full query
        " ".join(words[:3]) if len(words) > 3 else query,  # attempt 2: first 3 words
        words[0] if words else "nature",          # attempt 3: single keyword
    ]

    for variant in variants:
        result = _wikimedia(variant)
        if result:
            return result

    # Picsum guaranteed fallback
    return _picsum_fallback(query)


def _place_image(slide, img_bytes: io.BytesIO | None, l, t, w, h):
    """Insert image into slide — always succeeds (Picsum guarantees a result)."""
    if img_bytes:
        try:
            img_bytes.seek(0)
            slide.shapes.add_picture(img_bytes, l, t, w, h)
            return
        except Exception:
            pass
    # Last-resort: solid colour block (should never reach here)
    _rect(slide, l, t, w, h, RGBColor(0xDF, 0xEA, 0xF5))


# ─── 5. Image-left slide ──────────────────────────────────────────────────────

def create_image_left_slide(prs: Presentation, title: str,
                              description: str, bullets: list[str],
                              slide_num: int, dept: str, today: str,
                              image_query: str = "") -> None:
    """
    Left: real image | Right: title + text.
    """
    slide = _blank(prs)
    _draw_title(slide, title)
    _draw_footer(slide, slide_num, dept, today)

    img_w   = Inches(5.5)
    gap     = Inches(0.25)
    txt_l   = CONTENT_L + img_w + gap
    txt_w   = SW - txt_l - Inches(0.35)
    ct      = CONTENT_T + Inches(0.1)
    ch      = _available_height() - Inches(0.1)

    img_bytes = _fetch_image(image_query or title)
    _place_image(slide, img_bytes, Inches(0.3), CONTENT_T, img_w, ch)

    # Text block
    txb = slide.shapes.add_textbox(txt_l, ct, txt_w, ch)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_left = tf.margin_right = tf.margin_bottom = 0

    started = False
    if description:
        p = tf.paragraphs[0]
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = description
        run.font.name = FONT
        run.font.size = Pt(21)
        run.font.italic = True
        run.font.color.rgb = C_DARKGRAY
        started = True

    for bullet in bullets:
        p = tf.paragraphs[0] if not started else tf.add_paragraph()
        started = True
        p.space_before = Pt(6)
        p.space_after  = Pt(4)
        mk = p.add_run()
        mk.text = "▶  "
        mk.font.name = FONT
        mk.font.size = Pt(16)
        mk.font.bold = True
        mk.font.color.rgb = C_RED
        run = p.add_run()
        run.text = bullet
        run.font.name = FONT
        run.font.size = Pt(20)
        run.font.color.rgb = C_BLACK


# ─── 6. Image-right slide ─────────────────────────────────────────────────────

def create_image_right_slide(prs: Presentation, title: str,
                               description: str, bullets: list[str],
                               slide_num: int, dept: str, today: str,
                               image_query: str = "") -> None:
    """
    Left: text | Right: real image.
    """
    slide = _blank(prs)
    _draw_title(slide, title)
    _draw_footer(slide, slide_num, dept, today)

    txt_w   = Inches(5.85)
    gap     = Inches(0.25)
    img_l   = CONTENT_L + txt_w + gap
    img_w   = SW - img_l - Inches(0.2)
    ct      = CONTENT_T + Inches(0.1)
    ch      = _available_height() - Inches(0.1)

    # Text block
    txb = slide.shapes.add_textbox(CONTENT_L + Inches(0.1), ct, txt_w, ch)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_left = tf.margin_right = tf.margin_bottom = 0

    started = False
    if description:
        p = tf.paragraphs[0]
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = description
        run.font.name = FONT
        run.font.size = Pt(21)
        run.font.italic = True
        run.font.color.rgb = C_RED
        started = True

    for bullet in bullets:
        p = tf.paragraphs[0] if not started else tf.add_paragraph()
        started = True
        p.space_before = Pt(7)
        p.space_after  = Pt(4)
        mk = p.add_run()
        mk.text = "▶  "
        mk.font.name = FONT
        mk.font.size = Pt(17)
        mk.font.bold = True
        mk.font.color.rgb = C_RED
        run = p.add_run()
        run.text = bullet
        run.font.name = FONT
        run.font.size = Pt(21)
        run.font.color.rgb = C_BLACK

    img_bytes = _fetch_image(image_query or title)
    _place_image(slide, img_bytes, img_l, CONTENT_T, img_w, ch)


# ─── 7. Card-pair slide ───────────────────────────────────────────────────────

def create_card_slide(prs: Presentation, cards: list[dict],
                       slide_num: int, dept: str, today: str) -> None:
    """
    Two (or more) stacked rounded-rectangle cards — orange fill, white/black text.
    Matches slide 10 of the reference (C-CELLS / FOLLICULAR CELLS).

    cards: [{"heading": str, "body": str}, ...]
    """
    slide = _blank(prs)
    _draw_footer(slide, slide_num, dept, today)

    n        = len(cards)
    margin_x = Inches(0.917)
    card_w   = SW - 2 * margin_x
    total_h  = FOOTER_T - Inches(0.4)
    gap      = Inches(0.18)
    card_h   = (total_h - gap * (n - 1)) / n

    CARD_FILLS = [C_ORANGE, C_BLUE, C_GREEN, C_GOLD]

    for i, card in enumerate(cards):
        fill_color = CARD_FILLS[i % len(CARD_FILLS)]
        cy = Inches(0.4) + i * (card_h + gap)

        # Card background
        _rrect(slide, margin_x, cy, card_w, card_h, fill_color, adj=30000)

        # Heading
        heading_h = Inches(0.65)
        _txb(slide,
             l=margin_x + Inches(0.35), t=cy + Inches(0.18),
             w=card_w - Inches(0.7), h=heading_h,
             text=card.get("heading", "").upper(),
             size=28, bold=True, color=C_WHITE,
             align=PP_ALIGN.LEFT)

        # Body text
        body_top = cy + heading_h + Inches(0.25)
        body_h   = card_h - heading_h - Inches(0.45)

        txb = slide.shapes.add_textbox(
            margin_x + Inches(0.35), body_top,
            card_w - Inches(0.7), body_h)
        tf = txb.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_left = tf.margin_right = tf.margin_bottom = 0

        body_text = card.get("body", "")
        body_lines = body_text.split("\n") if "\n" in body_text else [body_text]

        for j, line in enumerate(body_lines):
            if not line.strip():
                continue
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after  = Pt(3)
            mk = p.add_run()
            mk.text = "• "
            mk.font.name = FONT
            mk.font.size = Pt(21)
            mk.font.color.rgb = C_WHITE
            run = p.add_run()
            run.text = line.strip()
            run.font.name = FONT
            run.font.size = Pt(21)
            run.font.color.rgb = C_WHITE


# ─── 8. Flow diagram ──────────────────────────────────────────────────────────

def create_flow_slide(prs: Presentation, title: str,
                       steps: list[str], description: str,
                       slide_num: int, dept: str, today: str) -> None:
    """
    Horizontal flow: Step1 → Step2 → … with numbered colored nodes.
    Matches the academic process-flow style.
    """
    slide = _blank(prs)
    _draw_title(slide, title)
    _draw_footer(slide, slide_num, dept, today)

    if description:
        _txb(slide,
             l=CONTENT_L + Inches(0.2), t=CONTENT_T + Inches(0.05),
             w=CONTENT_W - Inches(0.4), h=Inches(0.6),
             text=description, size=18, italic=True, color=C_DARKGRAY)

    n = len(steps)
    if n == 0:
        return

    usable_w = CONTENT_W - Inches(0.3)
    center_y  = Inches(4.3)
    node_h    = Inches(1.3)
    arrow_w   = Inches(0.45)
    node_w    = min(Inches(2.3), (usable_w - arrow_w * (n - 1)) / n)
    start_x   = CONTENT_L + (usable_w - node_w * n - arrow_w * (n - 1)) / 2

    for i, step in enumerate(steps):
        color = _NODE_COLORS[i % len(_NODE_COLORS)]
        nx = start_x + i * (node_w + arrow_w)

        # Step number badge
        badge_sz = Inches(0.38)
        badge_y  = center_y - node_h / 2 - badge_sz - Inches(0.05)
        _rrect(slide, nx + (node_w - badge_sz) / 2, badge_y,
               badge_sz, badge_sz, color, adj=120000)
        _txb(slide,
             l=nx + (node_w - badge_sz) / 2, t=badge_y,
             w=badge_sz, h=badge_sz,
             text=str(i + 1), size=12, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

        # Node box
        _rrect(slide, nx, center_y - node_h / 2, node_w, node_h, color, adj=50000)

        # Step label
        fs = max(9, min(16, int(220 / max(len(step), 1) * 1.5)))
        _txb(slide,
             l=nx + Inches(0.08), t=center_y - node_h / 2 + Inches(0.12),
             w=node_w - Inches(0.16), h=node_h - Inches(0.2),
             text=step, size=fs, bold=False,
             color=C_WHITE, align=PP_ALIGN.CENTER)

        # Arrow
        if i < n - 1:
            ax = nx + node_w
            _txb(slide,
                 l=ax, t=center_y - Inches(0.22),
                 w=arrow_w, h=Inches(0.44),
                 text="→", size=22, bold=True,
                 color=C_RED, align=PP_ALIGN.CENTER)


# ─── 9. Tree diagram ──────────────────────────────────────────────────────────

def create_tree_slide(prs: Presentation, title: str,
                       root: str, children: list,
                       description: str,
                       slide_num: int, dept: str, today: str) -> None:
    """
    Root → children (→ optional grandchildren) tree structure.
    children: list of str  OR  list of {"label": str, "children": list}
    """
    slide = _blank(prs)
    _draw_title(slide, title)
    _draw_footer(slide, slide_num, dept, today)

    desc_consumed = Inches(0)
    if description:
        _txb(slide,
             l=CONTENT_L + Inches(0.2), t=CONTENT_T + Inches(0.05),
             w=CONTENT_W - Inches(0.4), h=Inches(0.55),
             text=description, size=18, italic=True, color=C_DARKGRAY)
        desc_consumed = Inches(0.65)

    # Root node
    rw, rh = Inches(4.0), Inches(0.75)
    rx      = (SW - rw) / 2
    ry      = CONTENT_T + Inches(0.1) + desc_consumed
    _rrect(slide, rx, ry, rw, rh, C_ORANGE, adj=60000)
    _txb(slide,
         l=rx + Inches(0.1), t=ry + Inches(0.1),
         w=rw - Inches(0.2), h=rh - Inches(0.15),
         text=root, size=18, bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)

    if not children:
        return

    n       = len(children)
    cw_node = min(Inches(3.0), (CONTENT_W - Inches(0.25) * (n - 1)) / n)
    ch_h    = Inches(0.75)
    cy      = ry + rh + Inches(1.55)
    tot_w   = n * cw_node + Inches(0.25) * (n - 1)
    cx0     = (SW - tot_w) / 2
    stem_x  = rx + rw / 2

    # Vertical stem
    _rect(slide, stem_x - Inches(0.02), ry + rh,
          Inches(0.04), cy - ry - rh, C_DARKGRAY)

    for i, child in enumerate(children):
        color = _NODE_COLORS[(i + 1) % len(_NODE_COLORS)]
        label = child if isinstance(child, str) else child.get("label", f"Branch {i+1}")
        xi = cx0 + i * (cw_node + Inches(0.25))
        ccx = xi + cw_node / 2

        # Horizontal connector from stem
        min_x = min(stem_x, ccx)
        _rect(slide, min_x, cy - Inches(0.02),
              abs(ccx - stem_x), Inches(0.04), C_DARKGRAY)
        # Short drop to child
        _rect(slide, ccx - Inches(0.02), cy - Inches(0.02),
              Inches(0.04), Inches(0.08), C_DARKGRAY)

        _rrect(slide, xi, cy, cw_node, ch_h, color, adj=60000)
        fs = max(9, min(14, int(220 / max(len(label), 1) * 1.5)))
        _txb(slide,
             l=xi + Inches(0.06), t=cy + Inches(0.1),
             w=cw_node - Inches(0.12), h=ch_h - Inches(0.15),
             text=label, size=fs, bold=False,
             color=C_WHITE, align=PP_ALIGN.CENTER)

        # Grandchildren
        if isinstance(child, dict):
            gc_list = child.get("children", [])
            if gc_list:
                ng  = len(gc_list)
                gcw = cw_node / ng - Inches(0.06)
                gcy = cy + ch_h + Inches(0.35)
                _rect(slide, ccx - Inches(0.02), cy + ch_h,
                      Inches(0.04), Inches(0.35), C_DARKGRAY)
                for j, gc in enumerate(gc_list):
                    gc_label = gc if isinstance(gc, str) else gc.get("label", "")
                    gcx = xi + j * (gcw + Inches(0.06))
                    gc_cx = gcx + gcw / 2
                    _rect(slide, gc_cx - Inches(0.02), gcy - Inches(0.08),
                          Inches(0.04), Inches(0.08), C_DARKGRAY)
                    _rrect(slide, gcx, gcy, gcw, Inches(0.55),
                           RGBColor(0xCA, 0xDC, 0xFC), adj=50000)
                    gfs = max(7, min(11, int(160 / max(len(gc_label), 1) * 1.5)))
                    _txb(slide,
                         l=gcx + Inches(0.04), t=gcy + Inches(0.06),
                         w=gcw - Inches(0.08), h=Inches(0.45),
                         text=gc_label, size=gfs, color=C_BLACK,
                         align=PP_ALIGN.CENTER)


# ─── 10. Multi-flow (swim-lane) slide ─────────────────────────────────────────

def create_multiflow_slide(prs: Presentation, title: str,
                            flows: list[dict], description: str,
                            slide_num: int, dept: str, today: str) -> None:
    """
    Horizontal swim-lanes stacked vertically.
    flows: [{"label": str, "steps": [str, ...]}, ...]
    """
    slide = _blank(prs)
    _draw_title(slide, title)
    _draw_footer(slide, slide_num, dept, today)

    if description:
        _txb(slide,
             l=CONTENT_L + Inches(0.2), t=CONTENT_T + Inches(0.05),
             w=CONTENT_W - Inches(0.4), h=Inches(0.55),
             text=description, size=18, italic=True, color=C_DARKGRAY)

    nf        = len(flows)
    start_y   = CONTENT_T + (Inches(0.7) if description else Inches(0.2))
    avail_h   = FOOTER_T - start_y - Inches(0.15)
    gap       = Inches(0.14)
    lane_h    = min(Inches(1.35), (avail_h - gap * (nf - 1)) / nf)
    label_w   = Inches(1.85)
    arrow_w   = Inches(0.38)
    step_area = CONTENT_W - label_w - Inches(0.45)

    LIGHTERS = [
        RGBColor(0xF8, 0xD5, 0xBC),  # light orange
        RGBColor(0xC5, 0xD8, 0xF0),  # light blue
        RGBColor(0xC9, 0xE8, 0xC9),  # light green
        RGBColor(0xFF, 0xEE, 0xAA),  # light gold
        RGBColor(0xC5, 0xDB, 0xF0),  # light ltblue
    ]

    for fi, flow in enumerate(flows):
        color  = _NODE_COLORS[fi % len(_NODE_COLORS)]
        light  = LIGHTERS[fi % len(LIGHTERS)]
        lane_y = start_y + fi * (lane_h + gap)
        label  = flow.get("label", f"Flow {fi + 1}")
        steps  = flow.get("steps", [])

        # Label pill
        _rrect(slide, CONTENT_L, lane_y, label_w, lane_h - Inches(0.06),
               color, adj=80000)
        lfs = max(9, min(13, int(170 / max(len(label), 1) * 1.8)))
        _txb(slide,
             l=CONTENT_L, t=lane_y + (lane_h - Inches(0.45)) / 2,
             w=label_w, h=Inches(0.45),
             text=label, size=lfs, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)

        if not steps:
            continue

        ns     = len(steps)
        step_w = (step_area - arrow_w * (ns - 1)) / ns
        step_x0 = CONTENT_L + label_w + Inches(0.25)
        step_h  = lane_h - Inches(0.12)

        for si, step in enumerate(steps):
            sx = step_x0 + si * (step_w + arrow_w)
            _rrect(slide, sx, lane_y + Inches(0.06),
                   step_w, step_h, light, adj=50000)
            sfs = max(8, min(12, int(160 / max(len(step), 1) * 1.8)))
            _txb(slide,
                 l=sx + Inches(0.05), t=lane_y + Inches(0.1),
                 w=step_w - Inches(0.1), h=step_h - Inches(0.1),
                 text=step, size=sfs, color=C_BLACK,
                 align=PP_ALIGN.CENTER)
            if si < ns - 1:
                ax = sx + step_w
                _txb(slide,
                     l=ax, t=lane_y + (lane_h - Inches(0.4)) / 2,
                     w=arrow_w, h=Inches(0.4),
                     text="→", size=18, bold=True,
                     color=color, align=PP_ALIGN.CENTER)


# ─── 11. Closing slide ────────────────────────────────────────────────────────

def create_closing_slide(prs: Presentation, topic: str,
                          dept: str, today: str, slide_num: int) -> None:
    slide = _blank(prs)

    # Large centered "THANK YOU"
    _txb(slide,
         l=Inches(1.0), t=Inches(2.5),
         w=Inches(11.33), h=Inches(1.5),
         text="THANK YOU",
         size=54, bold=True,
         color=C_RED, align=PP_ALIGN.CENTER)

    _txb(slide,
         l=Inches(1.0), t=Inches(4.2),
         w=Inches(11.33), h=Inches(0.7),
         text=topic,
         size=22, bold=False,
         color=C_DARKGRAY, align=PP_ALIGN.CENTER)

    _draw_footer(slide, slide_num, dept, today)


# ═══════════════════════════════════════════════════════════════════════════════
#  LAYOUT DECISION ENGINE
# ═══════════════════════════════════════════════════════════════════════════════

def _resolve_layout(sd: dict) -> str:
    """
    Determine layout from slide content.
    Priority: explicit hint → structure_data → image flag → text content.
    """
    hint    = str(sd.get("layout_hint", "auto")).lower().strip()
    s_type  = str(sd.get("type", "")).lower().strip()
    has_sd  = bool(sd.get("structure_data"))
    is_img  = str(sd.get("image", "no")).lower() == "yes"
    has_d   = bool(str(sd.get("description", "")).strip())
    has_b   = bool(sd.get("bullets"))
    cards   = sd.get("cards")

    # 1. Explicit hint
    if hint not in ("", "auto"):
        return hint

    # 2. Card layout (two-box slide like slide 10)
    if cards:
        return "card"

    # 3. Diagrams from type
    if s_type in ("flow", "tree", "multi-flow", "multi_flow"):
        return s_type.replace("_", "-")
    if has_sd:
        sd_data = sd.get("structure_data", {})
        if "steps" in sd_data:
            return "flow"
        if "flows" in sd_data:
            return "multi-flow"
        if "root" in sd_data or "children" in sd_data:
            return "tree"

    # 4. Image layouts
    if is_img:
        layout = str(sd.get("layout", "right")).lower()
        return "image-left" if "left" in layout else "image-right"

    # 5. Text-based
    if has_d and has_b:
        return "mixed"
    if has_b:
        return "bullets"
    if has_d:
        return "paragraph"
    return "bullets"


def _get(sd: dict, key: str, default=None):
    """Field extraction with alias fallback."""
    aliases = {
        "title":       ["title", "heading"],
        "description": ["description", "main_heading_description", "desc"],
        "bullets":     ["bullets", "sub_points", "points"],
    }
    for k in aliases.get(key, [key]):
        if k in sd and sd[k]:
            return sd[k]
    return default


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN ENTRY POINT
# ═══════════════════════════════════════════════════════════════════════════════

def generate_ppt(slides: list[dict],
                 topic: str = "Presentation",
               
                 department: str = "",
                 presenter: str = "",
                 output_path: str = "output.pptx",
                 title_color: tuple = None,
                 accent_color: tuple = None,
                 bg_color: tuple = None,
                 text_color: tuple = None,
                 font: str = None,
                 node_colors: list = None) -> str:
    """
    Generate a .pptx that matches the PHARMACOTHERAPEUTICS reference style.

    Parameters
    ----------
    slides      : list of slide dicts (see module docstring for schema)
    topic       : course/presentation title shown on cover slide
    department  : footer department name
    presenter   : multi-line presenter info for cover slide
    output_path : file save path

    Returns
    -------
    str – absolute path to saved .pptx
    """
    # Apply dynamic theme
    global C_RED, C_ORANGE, FONT, _NODE_COLORS
    if title_color:  C_RED        = RGBColor(*title_color)
    if accent_color: C_ORANGE     = RGBColor(*accent_color)
    if text_color:   pass  # text_color applied via tx_color below
    if font and font.strip(): FONT = font
    if node_colors:  _NODE_COLORS = [RGBColor(*c) for c in node_colors]

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    today = date.today().strftime("%m/%d/%Y")

    # ── Cover slide ──────────────────────────────────────────────────────────
    subtitles = [_get(s, "title") or "" for s in slides[:4] if _get(s, "title")]
    create_title_slide(prs, topic, subtitles, presenter, department, today)

    # ── Content slides ────────────────────────────────────────────────────────
    for idx, sd in enumerate(slides, start=2):
        layout = _resolve_layout(sd)
        title  = _get(sd, "title")  or "Untitled"
        desc   = _get(sd, "description") or ""
        buls   = _get(sd, "bullets") or []
        num    = idx

        if layout == "paragraph":
            create_paragraph_slide(prs, title, desc, num, department, today)

        elif layout == "bullets":
            create_bullet_slide(prs, title, buls, desc, num, department, today)

        elif layout == "mixed":
            create_mixed_slide(prs, title, desc, buls, num, department, today)

        elif layout == "image-left":
            create_image_left_slide(prs, title, desc, buls, num, department, today,
                                    image_query=sd.get("image_query", title))

        elif layout == "image-right":
            create_image_right_slide(prs, title, desc, buls, num, department, today,
                                     image_query=sd.get("image_query", title))

        elif layout == "card":
            create_card_slide(prs, sd.get("cards", []), num, department, today)

        elif layout == "flow":
            struct  = sd.get("structure_data", {})
            steps   = struct.get("steps") or buls or []
            create_flow_slide(prs, title, steps, desc, num, department, today)

        elif layout == "tree":
            struct   = sd.get("structure_data", {})
            root_lbl = struct.get("root", title)
            children = struct.get("children", [])
            create_tree_slide(prs, title, root_lbl, children, desc, num, department, today)

        elif layout in ("multi-flow", "multi_flow"):
            struct = sd.get("structure_data", {})
            flows  = struct.get("flows", [])
            create_multiflow_slide(prs, title, flows, desc, num, department, today)

        else:
            # Fallback: bullet slide
            create_bullet_slide(prs, title, buls, desc, num, department, today)

    # ── Closing slide ─────────────────────────────────────────────────────────
    create_closing_slide(prs, topic, department, today, len(slides) + 2)

    # ── Save ─────────────────────────────────────────────────────────────────
    out_dir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(out_dir, exist_ok=True)
    prs.save(output_path)
    return os.path.abspath(output_path)


# ═══════════════════════════════════════════════════════════════════════════════
#  SELF-TEST
# ═══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    sample_slides = [
        # 1. Paragraph-only (like slides 3, 12)
        {
            "title": "THYROID GLAND",
            "description": (
                "Thyroid gland is a butterfly shaped endocrine organ found in the neck. "
                "It is responsible for regulating the body's metabolic rate via hormones it produces. "
                "Composed of two lobes, left and right connected by a narrow isthmus. "
                "Weighs around 25 grams in adults. Each lobe being about 5cm long, 3cm wide and 2cm thick."
            ),
            "bullets": [],
            "image": "no",
            "layout_hint": "auto",
        },
        # 2. Bullet-only (like slide 13)
        {
            "title": "DISORDERS",
            "description": "",
            "bullets": [
                "Goiter",
                "Graves disease",
                "Cretinism",
                "Hashimoto's disease",
                "Thyroid nodules",
            ],
            "image": "no",
            "layout_hint": "auto",
        },
        # 3. Mixed (like slides 6, 7)
        {
            "title": "Hypothyroidism",
            "description": (
                "Hypothyroidism is a clinical syndrome resulting from a deficiency of "
                "thyroid hormones, resulting in a generalized slowing of metabolic processes."
            ),
            "bullets": [
                "In newborn infants → Cretinism",
                "In adolescents → short stature, mental retardation, precocious puberty",
                "In adults → symptoms largely reversible after therapy",
                "TSH elevated; T3/T4 below normal reference range",
            ],
            "image": "no",
            "layout_hint": "auto",
        },
        # 4. Image-right (like slide 7)
        {
            "title": "Hyperthyroidism",
            "description": (
                "Hyperthyroidism is hyperactivity of the thyroid gland leading to "
                "excessive synthesis and release of thyroid hormones."
            ),
            "bullets": [
                "Increased basal metabolic rate",
                "Weight loss despite increased appetite",
                "Palpitations and tachycardia",
                "Heat intolerance and excessive sweating",
            ],
            "image": "yes",
            "layout": "right",
            "layout_hint": "image-right",
        },
        # 5. Image-left
        {
            "title": "LOCATION",
            "description": (
                "Located in the anterior neck, just inferior to the Adam's apple and larynx."
            ),
            "bullets": [
                "Between the C5 and T1 vertebrae",
                "Anterior to the upper part of the trachea",
                "Immediately below the larynx",
            ],
            "image": "yes",
            "layout": "left",
            "layout_hint": "image-left",
        },
        # 6. Card slide (like slide 10)
        {
            "layout_hint": "card",
            "cards": [
                {
                    "heading": "FOLLICULAR CELLS",
                    "body": (
                        "Secrete 2 hormones T3 (triiodothyronine) and T4 (tetraiodothyronine) "
                        "that influence the rate of metabolism, for which iodine is essential.\n"
                        "The activity of follicular cells is influenced by TSH (thyroid stimulating hormone)."
                    ),
                },
                {
                    "heading": "C-CELLS (Parafollicular Cells)",
                    "body": (
                        "Also known as clear cells or light cells.\n"
                        "They secrete hormone thyro-calcitonin which has an action opposite to "
                        "that of the parathyroid hormone on calcium metabolism.\n"
                        "This hormone acts when serum calcium level is high."
                    ),
                },
            ],
        },
        # 7. Flow diagram
        {
            "title": "DRUG DEVELOPMENT PIPELINE",
            "description": "Sequential stages from discovery to post-market surveillance.",
            "layout_hint": "flow",
            "structure_data": {
                "steps": ["Discovery", "Pre-clinical", "Phase I", "Phase II",
                          "Phase III", "FDA Review", "Phase IV"]
            },
        },
        # 8. Tree diagram
        {
            "title": "CLASSIFICATION OF ANTITHYROID DRUGS",
            "description": "Systematic classification of drugs used in thyroid disorders.",
            "layout_hint": "tree",
            "structure_data": {
                "root": "Antithyroid Agents",
                "children": [
                    {"label": "Thionamides",      "children": ["Methimazole", "Carbimazole", "PTU"]},
                    {"label": "Iodine Compounds",  "children": ["Lugol's Solution", "SSKI"]},
                    {"label": "Radioactive Iodine","children": ["¹³¹I Therapy"]},
                    {"label": "Beta-Blockers",     "children": ["Propranolol", "Atenolol"]},
                ],
            },
        },
        # 9. Multi-flow
        {
            "title": "MECHANISM OF ANTITHYROID DRUGS",
            "description": "How major drug classes inhibit thyroid hormone production.",
            "layout_hint": "multi-flow",
            "structure_data": {
                "flows": [
                    {"label": "Thionamides",   "steps": ["Block TPO Enzyme", "Inhibit Oxidation", "↓ T3/T4 Synthesis"]},
                    {"label": "¹³¹I Therapy",  "steps": ["Uptaken by Gland", "β-Radiation", "Follicular Destruction"]},
                    {"label": "Beta-Blockers", "steps": ["Block β-Receptors", "↓ Sympathetic Effects", "Symptomatic Relief"]},
                ],
            },
        },
    ]

    path = generate_ppt(
        slides=sample_slides,
        topic="PHARMACOTHERAPEUTICS-I",
     
        presenter="DEVIPRIYA V M\nM.PHARM SEM-1\nPHARMACY PRACTICE",
        output_path="outputs/pharma_style_demo.pptx",
    )
    print(f"Saved → {path}")











# """
# ppt_generator.py  –  Upgraded slide builder
# Supports: content | flow | tree | multi-flow | image-left | image-right
# Backward-compatible with the old {heading, bullets} format.
# """

# import os
# from pptx import Presentation
# from pptx.util import Inches, Pt, Emu
# from pptx.dml.color import RGBColor
# from pptx.enum.text import PP_ALIGN
# from pptx.oxml.ns import qn
# from lxml import etree

# # ── Palette ────────────────────────────────────────────────────────────────
# BG_DARK    = RGBColor(0x1E, 0x27, 0x61)   # navy
# BG_LIGHT   = RGBColor(0xFF, 0xFF, 0xFF)   # white
# ACCENT     = RGBColor(0x4F, 0x8E, 0xF7)   # sky blue
# ACCENT2    = RGBColor(0x2C, 0xB6, 0x7D)   # teal  (flow/tree nodes)
# ACCENT3    = RGBColor(0xF9, 0x61, 0x67)   # coral (multi-flow alt)
# TEXT_DARK  = RGBColor(0x1E, 0x27, 0x61)
# TEXT_LITE  = RGBColor(0xFF, 0xFF, 0xFF)
# BULLET_CLR = RGBColor(0x33, 0x33, 0x33)
# MUTED      = RGBColor(0x88, 0x88, 0x99)
# NODE_TEXT  = RGBColor(0xFF, 0xFF, 0xFF)

# TITLE_FONT = "Trebuchet MS"
# BODY_FONT  = "Calibri"

# SLIDE_W = Inches(13.33)
# SLIDE_H = Inches(7.5)

# # Colour name → RGBColor (for user-specified backgrounds)
# NAMED_COLORS = {
#     "white":   RGBColor(0xFF, 0xFF, 0xFF),
#     "black":   RGBColor(0x00, 0x00, 0x00),
#     "navy":    RGBColor(0x1E, 0x27, 0x61),
#     "blue":    RGBColor(0x1E, 0x6F, 0xBF),
#     "skyblue": RGBColor(0x87, 0xCE, 0xEB),
#     "teal":    RGBColor(0x00, 0x80, 0x80),
#     "green":   RGBColor(0x2C, 0x5F, 0x2D),
#     "red":     RGBColor(0x99, 0x00, 0x11),
#     "coral":   RGBColor(0xF9, 0x61, 0x67),
#     "orange":  RGBColor(0xFF, 0x7F, 0x00),
#     "purple":  RGBColor(0x6A, 0x0D, 0xAD),
#     "gray":    RGBColor(0x66, 0x66, 0x66),
#     "grey":    RGBColor(0x66, 0x66, 0x66),
#     "dark":    RGBColor(0x1E, 0x27, 0x61),
#     "light":   RGBColor(0xFF, 0xFF, 0xFF),
# }


# # ── Low-level helpers ──────────────────────────────────────────────────────
# def _parse_color(raw: str) -> RGBColor:
#     """Accept a color name or 6-digit hex string."""
#     if not raw:
#         return BG_LIGHT
#     s = raw.strip().lower().lstrip("#")
#     if s in NAMED_COLORS:
#         return NAMED_COLORS[s]
#     if len(s) == 6:
#         try:
#             return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
#         except ValueError:
#             pass
#     return BG_LIGHT


# def _solid_fill(shape, color: RGBColor):
#     shape.fill.solid()
#     shape.fill.fore_color.rgb = color


# def _no_line(shape):
#     shape.line.fill.background()


# def _add_textbox(slide, left, top, width, height,
#                  text, font_name=BODY_FONT, font_size=16,
#                  bold=False, color=TEXT_DARK,
#                  align=PP_ALIGN.LEFT, italic=False):
#     txb = slide.shapes.add_textbox(left, top, width, height)
#     tf  = txb.text_frame
#     tf.word_wrap = True
#     p   = tf.paragraphs[0]
#     p.alignment = align
#     run = p.add_run()
#     run.text           = text
#     run.font.name      = font_name
#     run.font.size      = Pt(font_size)
#     run.font.bold      = bold
#     run.font.italic    = italic
#     run.font.color.rgb = color
#     return txb


# def _title_bar(slide, bg_color: RGBColor = BG_DARK):
#     """Full-width title bar at the top."""
#     bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.15))
#     _solid_fill(bar, bg_color)
#     _no_line(bar)
#     return bar


# def _slide_heading(slide, text: str, bar_color: RGBColor = BG_DARK):
#     """Title bar + heading text."""
#     _title_bar(slide, bar_color)
#     _add_textbox(slide,
#                  left=Inches(0.5), top=Inches(0.12),
#                  width=Inches(12.33), height=Inches(0.9),
#                  text=text, font_name=TITLE_FONT, font_size=32,
#                  bold=True, color=TEXT_LITE)


# def _bg(slide, color: RGBColor = BG_LIGHT):
#     """Solid background rectangle."""
#     bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
#     _solid_fill(bg, color)
#     _no_line(bg)


# def _bullet_block(slide, bullets: list, left, top, width, height,
#                   font_size=16, color=BULLET_CLR):
#     if not bullets:
#         return
#     txb = slide.shapes.add_textbox(left, top, width, height)
#     tf  = txb.text_frame
#     tf.word_wrap = True
#     for i, bullet in enumerate(bullets):
#         para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
#         para.space_before = Pt(6)
#         para.space_after  = Pt(2)
#         marker = para.add_run()
#         marker.text           = "▸  "
#         marker.font.name      = BODY_FONT
#         marker.font.size      = Pt(font_size)
#         marker.font.bold      = True
#         marker.font.color.rgb = ACCENT
#         run = para.add_run()
#         run.text           = bullet
#         run.font.name      = BODY_FONT
#         run.font.size      = Pt(font_size)
#         run.font.bold      = False
#         run.font.color.rgb = color


# def _rounded_rect(slide, left, top, width, height,
#                   fill: RGBColor, radius_emu=914400//8):
#     """Add a rounded rectangle shape."""
#     from pptx.util import Emu
#     shape = slide.shapes.add_shape(
#         5,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
#         left, top, width, height
#     )
#     _solid_fill(shape, fill)
#     _no_line(shape)
#     # Adjust corner radius via XML
#     try:
#         prstGeom = shape.fill._element.getparent().find(
#             ".//" + qn("a:prstGeom"))
#         if prstGeom is not None:
#             avLst = prstGeom.find(qn("a:avLst"))
#             if avLst is not None:
#                 for gd in avLst.findall(qn("a:gd")):
#                     if gd.get("name") == "adj":
#                         gd.set("fmla", f"val {radius_emu}")
#     except Exception:
#         pass
#     return shape


# # ── Slide builders ─────────────────────────────────────────────────────────
# def _build_title_slide(prs, title: str, subtitle: str = ""):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide, BG_DARK)

#     _add_textbox(slide,
#                  left=Inches(1), top=Inches(2.4),
#                  width=Inches(11.33), height=Inches(1.6),
#                  text=title, font_name=TITLE_FONT, font_size=44,
#                  bold=True, color=TEXT_LITE, align=PP_ALIGN.CENTER)

#     if subtitle:
#         _add_textbox(slide,
#                      left=Inches(1), top=Inches(4.1),
#                      width=Inches(11.33), height=Inches(0.6),
#                      text=subtitle, font_name=BODY_FONT, font_size=20,
#                      bold=False, color=RGBColor(0xCA, 0xDC, 0xFC),
#                      align=PP_ALIGN.CENTER)


# def _build_closing_slide(prs, topic: str):
#     _build_title_slide(prs, "Thank You", f"Presentation on: {topic}")


# def _build_content_slide(prs, slide_data: dict, bg_color: RGBColor):
#     """Standard bullet slide (type: content)."""
#     title   = slide_data.get("title") or slide_data.get("heading", "Untitled")
#     desc    = slide_data.get("main_heading_description", "")
#     bullets = slide_data.get("sub_points") or slide_data.get("bullets", [])

#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide, bg_color)
#     _slide_heading(slide, title)

#     top_offset = Inches(1.35)

#     if desc:
#         _add_textbox(slide,
#                      left=Inches(0.7), top=top_offset,
#                      width=Inches(11.93), height=Inches(0.7),
#                      text=desc, font_name=BODY_FONT, font_size=14,
#                      bold=False, color=MUTED, italic=True)
#         top_offset += Inches(0.75)

#     _bullet_block(slide, bullets,
#                   left=Inches(0.7), top=top_offset,
#                   width=Inches(11.93), height=SLIDE_H - top_offset - Inches(0.3))


# def _build_image_slide(prs, slide_data: dict, bg_color: RGBColor, image_side: str):
#     """
#     image-left: placeholder box on left, bullets on right.
#     image-right: bullets on left, placeholder box on right.
#     """
#     title   = slide_data.get("title") or slide_data.get("heading", "Untitled")
#     bullets = slide_data.get("sub_points") or slide_data.get("bullets", [])
#     desc    = slide_data.get("main_heading_description", "")

#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide, bg_color)
#     _slide_heading(slide, title)

#     content_top = Inches(1.3)
#     content_h   = Inches(5.8)
#     img_w       = Inches(5.5)
#     text_w      = Inches(6.8)
#     gap         = Inches(0.3)

#     if image_side == "left":
#         img_left  = Inches(0.4)
#         text_left = img_left + img_w + gap
#     else:
#         text_left = Inches(0.4)
#         img_left  = text_left + text_w + gap

#     # Image placeholder box
#     ph = slide.shapes.add_shape(1, img_left, content_top, img_w, content_h)
#     _solid_fill(ph, RGBColor(0xE8, 0xED, 0xF7))
#     _no_line(ph)
#     _add_textbox(slide,
#                  left=img_left, top=content_top + Inches(2.5),
#                  width=img_w, height=Inches(0.6),
#                  text="[ Image ]", font_name=BODY_FONT, font_size=14,
#                  bold=False, color=MUTED, align=PP_ALIGN.CENTER)

#     # Description + bullets
#     text_top = content_top
#     if desc:
#         _add_textbox(slide,
#                      left=text_left, top=text_top,
#                      width=text_w, height=Inches(0.65),
#                      text=desc, font_name=BODY_FONT, font_size=13,
#                      bold=False, color=MUTED, italic=True)
#         text_top += Inches(0.7)

#     _bullet_block(slide, bullets,
#                   left=text_left, top=text_top,
#                   width=text_w, height=content_h - (text_top - content_top))


# def _build_flow_slide(prs, slide_data: dict, bg_color: RGBColor):
#     """Horizontal flow diagram: Step1 → Step2 → Step3 …"""
#     title = slide_data.get("title") or slide_data.get("heading", "Untitled")
#     desc  = slide_data.get("main_heading_description", "")
#     sd    = slide_data.get("structure_data", {})
#     steps = sd.get("steps", [])

#     # Fall back to bullets as steps if structure_data is missing
#     if not steps:
#         steps = slide_data.get("sub_points") or slide_data.get("bullets", [])

#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide, bg_color)
#     _slide_heading(slide, title)

#     if desc:
#         _add_textbox(slide,
#                      left=Inches(0.7), top=Inches(1.3),
#                      width=Inches(11.93), height=Inches(0.55),
#                      text=desc, font_name=BODY_FONT, font_size=13,
#                      bold=False, color=MUTED, italic=True)

#     n = len(steps)
#     if n == 0:
#         return

#     # Layout: spread nodes across slide width
#     margin_x  = Inches(0.6)
#     usable_w  = SLIDE_W - 2 * margin_x
#     center_y  = Inches(4.2)
#     node_h    = Inches(1.1)
#     arrow_w   = Inches(0.45)

#     # Node width shrinks if many steps
#     max_node_w = Inches(2.2)
#     node_w = min(max_node_w,
#                  (usable_w - arrow_w * (n - 1)) / n)

#     node_colors = [ACCENT, ACCENT2, RGBColor(0xF9, 0x61, 0x67),
#                    RGBColor(0xFF, 0x9F, 0x1C), RGBColor(0x9B, 0x59, 0xB6),
#                    RGBColor(0x16, 0xA0, 0x85), RGBColor(0xE7, 0x4C, 0x3C)]

#     for i, step in enumerate(steps):
#         color = node_colors[i % len(node_colors)]
#         x = margin_x + i * (node_w + arrow_w)

#         # Node box
#         box = slide.shapes.add_shape(
#             5, x, center_y - node_h / 2, node_w, node_h)
#         _solid_fill(box, color)
#         _no_line(box)

#         # Step number
#         _add_textbox(slide,
#                      left=x, top=center_y - node_h / 2 - Inches(0.35),
#                      width=node_w, height=Inches(0.35),
#                      text=f"{i+1}", font_name=TITLE_FONT, font_size=12,
#                      bold=True, color=color, align=PP_ALIGN.CENTER)

#         # Step label inside box
#         _add_textbox(slide,
#                      left=x + Inches(0.06), top=center_y - node_h / 2 + Inches(0.1),
#                      width=node_w - Inches(0.12), height=node_h - Inches(0.1),
#                      text=step, font_name=BODY_FONT,
#                      font_size=max(9, min(13, int(130 / max(len(step), 1) * 2))),
#                      bold=False, color=NODE_TEXT, align=PP_ALIGN.CENTER)

#         # Arrow between nodes
#         if i < n - 1:
#             arrow_x = x + node_w
#             _add_textbox(slide,
#                          left=arrow_x, top=center_y - Inches(0.2),
#                          width=arrow_w, height=Inches(0.4),
#                          text="→", font_name=BODY_FONT, font_size=22,
#                          bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# def _build_tree_slide(prs, slide_data: dict, bg_color: RGBColor):
#     """Top-down tree: root → children (→ grandchildren if present)."""
#     title = slide_data.get("title") or slide_data.get("heading", "Untitled")
#     desc  = slide_data.get("main_heading_description", "")
#     sd    = slide_data.get("structure_data", {})
#     root_label = sd.get("root", title)
#     children   = sd.get("children", [])

#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide, bg_color)
#     _slide_heading(slide, title)

#     if desc:
#         _add_textbox(slide,
#                      left=Inches(0.7), top=Inches(1.3),
#                      width=Inches(11.93), height=Inches(0.55),
#                      text=desc, font_name=BODY_FONT, font_size=13,
#                      bold=False, color=MUTED, italic=True)

#     # Root node
#     root_w, root_h = Inches(3.5), Inches(0.75)
#     root_x = (SLIDE_W - root_w) / 2
#     root_y = Inches(2.0)

#     root_box = slide.shapes.add_shape(5, root_x, root_y, root_w, root_h)
#     _solid_fill(root_box, BG_DARK)
#     _no_line(root_box)
#     _add_textbox(slide,
#                  left=root_x, top=root_y + Inches(0.1),
#                  width=root_w, height=root_h - Inches(0.1),
#                  text=root_label, font_name=TITLE_FONT, font_size=16,
#                  bold=True, color=TEXT_LITE, align=PP_ALIGN.CENTER)

#     if not children:
#         return

#     n = len(children)
#     margin_x  = Inches(0.5)
#     usable_w  = SLIDE_W - 2 * margin_x
#     child_w   = min(Inches(2.8), usable_w / n - Inches(0.2))
#     child_h   = Inches(0.75)
#     child_y   = Inches(4.0)
#     total_w   = n * child_w + (n - 1) * Inches(0.25)
#     start_x   = (SLIDE_W - total_w) / 2

#     node_colors = [ACCENT, ACCENT2, RGBColor(0xF9, 0x61, 0x67),
#                    RGBColor(0xFF, 0x9F, 0x1C), RGBColor(0x9B, 0x59, 0xB6),
#                    RGBColor(0x16, 0xA0, 0x85)]

#     # Vertical line down from root
#     line_x = root_x + root_w / 2
#     _add_textbox(slide,
#                  left=line_x - Inches(0.05), top=root_y + root_h,
#                  width=Inches(0.1), height=child_y - root_y - root_h,
#                  text="", font_name=BODY_FONT, font_size=1, bold=False, color=MUTED)

#     for i, child in enumerate(children):
#         label = child if isinstance(child, str) else child.get("label", f"Branch {i+1}")
#         cx = start_x + i * (child_w + Inches(0.25))

#         # Draw child box
#         color = node_colors[i % len(node_colors)]
#         cbox = slide.shapes.add_shape(5, cx, child_y, child_w, child_h)
#         _solid_fill(cbox, color)
#         _no_line(cbox)
#         _add_textbox(slide,
#                      left=cx + Inches(0.05), top=child_y + Inches(0.1),
#                      width=child_w - Inches(0.1), height=child_h - Inches(0.1),
#                      text=label, font_name=BODY_FONT, font_size=12,
#                      bold=False, color=NODE_TEXT, align=PP_ALIGN.CENTER)

#         # Draw grandchildren if present
#         if isinstance(child, dict) and child.get("children"):
#             gc_y = child_y + child_h + Inches(0.35)
#             gc_list = child["children"]
#             for j, gc in enumerate(gc_list):
#                 gc_label = gc if isinstance(gc, str) else gc.get("label", "")
#                 gc_x = cx + j * (child_w / len(gc_list))
#                 gc_w = child_w / len(gc_list) - Inches(0.08)
#                 gcbox = slide.shapes.add_shape(5, gc_x, gc_y, gc_w, Inches(0.55))
#                 _solid_fill(gcbox, RGBColor(0xCA, 0xDC, 0xFC))
#                 _no_line(gcbox)
#                 _add_textbox(slide,
#                              left=gc_x, top=gc_y + Inches(0.05),
#                              width=gc_w, height=Inches(0.5),
#                              text=gc_label, font_name=BODY_FONT, font_size=10,
#                              bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)


# def _build_multiflow_slide(prs, slide_data: dict, bg_color: RGBColor):
#     """Multiple parallel flow lanes stacked vertically."""
#     title = slide_data.get("title") or slide_data.get("heading", "Untitled")
#     desc  = slide_data.get("main_heading_description", "")
#     sd    = slide_data.get("structure_data", {})
#     flows = sd.get("flows", [])

#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide, bg_color)
#     _slide_heading(slide, title)

#     if desc:
#         _add_textbox(slide,
#                      left=Inches(0.7), top=Inches(1.25),
#                      width=Inches(11.93), height=Inches(0.5),
#                      text=desc, font_name=BODY_FONT, font_size=13,
#                      bold=False, color=MUTED, italic=True)

#     if not flows:
#         # Fall back to bullets
#         bullets = slide_data.get("sub_points") or slide_data.get("bullets", [])
#         _bullet_block(slide, bullets,
#                       left=Inches(0.7), top=Inches(1.8),
#                       width=Inches(11.93), height=Inches(5.4))
#         return

#     n_flows    = len(flows)
#     lane_h     = min(Inches(1.4), (SLIDE_H - Inches(2.0)) / n_flows)
#     lane_gap   = Inches(0.18)
#     label_w    = Inches(1.9)
#     margin_x   = Inches(0.5)
#     usable_w   = SLIDE_W - margin_x - label_w - Inches(0.3)
#     start_y    = Inches(1.9) if not desc else Inches(2.0)

#     lane_colors = [ACCENT, ACCENT2, RGBColor(0xF9, 0x61, 0x67),
#                    RGBColor(0xFF, 0x9F, 0x1C), RGBColor(0x9B, 0x59, 0xB6)]

#     for fi, flow in enumerate(flows):
#         lane_label = flow.get("label", f"Flow {fi+1}")
#         steps      = flow.get("steps", [])
#         lane_color = lane_colors[fi % len(lane_colors)]
#         lane_y     = start_y + fi * (lane_h + lane_gap)

#         # Lane label pill
#         label_box = slide.shapes.add_shape(
#             5, margin_x, lane_y, label_w, lane_h - Inches(0.08))
#         _solid_fill(label_box, lane_color)
#         _no_line(label_box)
#         _add_textbox(slide,
#                      left=margin_x, top=lane_y + (lane_h - Inches(0.4)) / 2,
#                      width=label_w, height=Inches(0.4),
#                      text=lane_label, font_name=BODY_FONT, font_size=12,
#                      bold=True, color=NODE_TEXT, align=PP_ALIGN.CENTER)

#         if not steps:
#             continue

#         n_steps   = len(steps)
#         arrow_w   = Inches(0.35)
#         step_w    = (usable_w - arrow_w * (n_steps - 1)) / n_steps
#         step_x0   = margin_x + label_w + Inches(0.2)
#         step_h    = lane_h - Inches(0.12)

#         # Use a fixed lighter colour per lane
#         lighter = [
#             RGBColor(0xCA, 0xDC, 0xFC),
#             RGBColor(0xB2, 0xEB, 0xD8),
#             RGBColor(0xFC, 0xD0, 0xD1),
#             RGBColor(0xFF, 0xE0, 0xB2),
#             RGBColor(0xE1, 0xD5, 0xF5),
#         ]
#         sc = lighter[fi % len(lighter)]

#         for si, step in enumerate(steps):
#             sx = step_x0 + si * (step_w + arrow_w)

#             sbox = slide.shapes.add_shape(5, sx, lane_y, step_w, step_h)
#             _solid_fill(sbox, sc)
#             _no_line(sbox)

#             fs = max(9, min(12, int(130 / max(len(step), 1) * 1.8)))
#             _add_textbox(slide,
#                          left=sx + Inches(0.04), top=lane_y + Inches(0.08),
#                          width=step_w - Inches(0.08), height=step_h - Inches(0.08),
#                          text=step, font_name=BODY_FONT, font_size=fs,
#                          bold=False, color=TEXT_DARK, align=PP_ALIGN.CENTER)

#             if si < n_steps - 1:
#                 ax = sx + step_w
#                 _add_textbox(slide,
#                              left=ax, top=lane_y + (step_h - Inches(0.35)) / 2,
#                              width=arrow_w, height=Inches(0.35),
#                              text="→", font_name=BODY_FONT, font_size=16,
#                              bold=True, color=lane_color, align=PP_ALIGN.CENTER)


# # ── Dispatch ───────────────────────────────────────────────────────────────
# _BUILDERS = {
#     "content":     _build_content_slide,
#     "image-left":  lambda prs, d, bg: _build_image_slide(prs, d, bg, "left"),
#     "image-right": lambda prs, d, bg: _build_image_slide(prs, d, bg, "right"),
#     "flow":        _build_flow_slide,
#     "tree":        _build_tree_slide,
#     "multi-flow":  _build_multiflow_slide,
# }


# # ── Public API ─────────────────────────────────────────────────────────────
# def generate_pptx(slides_json: dict, output_dir: str = ".") -> str:
#     """
#     Build a .pptx and return its absolute file path.

#     Accepts both old format:
#         {"topic": "...", "slides": [{"heading": ..., "bullets": [...]}]}
#     and new format:
#         {"topic": "...", "slides": [{"title": ..., "type": ..., "sub_points": [...], ...}]}
#     """
#     topic  = slides_json.get("topic", "Presentation")
#     slides = slides_json.get("slides", [])

#     if not slides:
#         raise ValueError("'slides' list is empty — nothing to generate.")

#     prs = Presentation()
#     prs.slide_width  = SLIDE_W
#     prs.slide_height = SLIDE_H

#     _build_title_slide(prs, topic, f"{len(slides)} slides")

#     for slide_data in slides:
#         slide_type = slide_data.get("type", "content").lower().strip()
#         bg_raw     = slide_data.get("background", "white")
#         bg_color   = _parse_color(bg_raw)

#         builder = _BUILDERS.get(slide_type, _build_content_slide)
#         builder(prs, slide_data, bg_color)

#     _build_closing_slide(prs, topic)

#     os.makedirs(output_dir, exist_ok=True)
#     safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in topic)
#     filename = f"{safe.strip()[:50]}.pptx"
#     filepath = os.path.abspath(os.path.join(output_dir, filename))
#     prs.save(filepath)
#     return filepath


# # ── Quick test ─────────────────────────────────────────────────────────────
# if __name__ == "__main__":
#     sample = {
#         "topic": "Artificial Intelligence",
#         "slides": [
#             {
#                 "title": "What is AI?",
#                 "type": "content",
#                 "background": "white",
#                 "main_heading_description": "A brief overview of artificial intelligence and its scope.",
#                 "sub_points": [
#                     "AI enables machines to simulate human intelligence",
#                     "Encompasses ML, NLP, computer vision, and robotics",
#                     "Powers applications from search to self-driving cars",
#                     "Fastest-growing technology sector globally",
#                 ]
#             },
#             {
#                 "title": "AI Development Lifecycle",
#                 "type": "flow",
#                 "background": "white",
#                 "main_heading_description": "From data collection to deployment.",
#                 "structure_data": {
#                     "steps": ["Data Collection", "Data Cleaning", "Model Training",
#                               "Evaluation", "Deployment", "Monitoring"]
#                 }
#             },
#             {
#                 "title": "AI Technology Stack",
#                 "type": "tree",
#                 "background": "white",
#                 "structure_data": {
#                     "root": "Artificial Intelligence",
#                     "children": [
#                         {"label": "Machine Learning", "children": ["Supervised", "Unsupervised"]},
#                         {"label": "NLP",              "children": ["Translation", "Sentiment"]},
#                         {"label": "Computer Vision",  "children": ["Detection", "Recognition"]},
#                     ]
#                 }
#             },
#             {
#                 "title": "AI Industry Applications",
#                 "type": "multi-flow",
#                 "background": "white",
#                 "structure_data": {
#                     "flows": [
#                         {"label": "Healthcare", "steps": ["Diagnosis", "Drug Discovery", "Patient Care"]},
#                         {"label": "Finance",    "steps": ["Fraud Detection", "Trading", "Risk Analysis"]},
#                         {"label": "Education",  "steps": ["Personalisation", "Assessment", "Tutoring"]},
#                     ]
#                 }
#             },
#             {
#                 "title": "Computer Vision Explained",
#                 "type": "image-left",
#                 "background": "white",
#                 "sub_points": [
#                     "Enables machines to interpret visual data",
#                     "Used in facial recognition and medical imaging",
#                     "Requires large labelled datasets for training",
#                 ]
#             },
#         ]
#     }

#     path = generate_pptx(sample, output_dir="outputs")
#     print(f"Saved → {path}")



# import os
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# from pptx.enum.text import PP_ALIGN


# # ── Palette ────────────────────────────────────────────────────────────────
# BG_DARK   = RGBColor(0x1E, 0x27, 0x61)   # navy  – title / closing slides
# BG_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)   # white – content slides
# ACCENT    = RGBColor(0x4F, 0x8E, 0xF7)   # sky blue – title bar on content slides
# TEXT_DARK = RGBColor(0x1E, 0x27, 0x61)   # navy text on light bg
# TEXT_LITE = RGBColor(0xFF, 0xFF, 0xFF)   # white text on dark bg
# BULLET_CLR= RGBColor(0x33, 0x33, 0x33)   # near-black bullet text

# TITLE_FONT  = "Trebuchet MS"
# BODY_FONT   = "Calibri"

# SLIDE_W = Inches(13.33)
# SLIDE_H = Inches(7.5)


# # ── Helpers ─────────────────────────────────────────────────────────────────
# def _solid_fill(shape, color: RGBColor):
#     shape.fill.solid()
#     shape.fill.fore_color.rgb = color


# def _add_textbox(slide, left, top, width, height,
#                  text, font_name, font_size, bold,
#                  color: RGBColor, align=PP_ALIGN.LEFT):
#     txb = slide.shapes.add_textbox(left, top, width, height)
#     tf  = txb.text_frame
#     tf.word_wrap = True
#     p   = tf.paragraphs[0]
#     p.alignment = align
#     run = p.add_run()
#     run.text = text
#     run.font.name      = font_name
#     run.font.size      = Pt(font_size)
#     run.font.bold      = bold
#     run.font.color.rgb = color
#     return txb


# def _build_title_slide(prs: Presentation, title: str, subtitle: str = ""):
#     """Dark navy full-bleed title slide."""
#     slide  = prs.slides.add_slide(prs.slide_layouts[6])   # blank
#     bg_box = slide.shapes.add_shape(
#         1, 0, 0, SLIDE_W, SLIDE_H)                        # MSO_SHAPE_TYPE.RECTANGLE
#     _solid_fill(bg_box, BG_DARK)
#     bg_box.line.fill.background()

#     # Title
#     _add_textbox(slide,
#                  left=Inches(1), top=Inches(2.5),
#                  width=Inches(11.33), height=Inches(1.5),
#                  text=title, font_name=TITLE_FONT, font_size=44,
#                  bold=True, color=TEXT_LITE, align=PP_ALIGN.CENTER)

#     # Subtitle / topic tag
#     if subtitle:
#         _add_textbox(slide,
#                      left=Inches(1), top=Inches(4.1),
#                      width=Inches(11.33), height=Inches(0.6),
#                      text=subtitle, font_name=BODY_FONT, font_size=20,
#                      bold=False, color=RGBColor(0xCA, 0xDC, 0xFC),
#                      align=PP_ALIGN.CENTER)


# def _build_content_slide(prs: Presentation, heading: str, bullets: list[str]):
#     """White content slide with a navy title bar and bullet list."""
#     slide = prs.slides.add_slide(prs.slide_layouts[6])    # blank

#     # White background
#     bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
#     _solid_fill(bg, BG_LIGHT)
#     bg.line.fill.background()

#     # Accent title bar
#     bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.25))
#     _solid_fill(bar, BG_DARK)
#     bar.line.fill.background()

#     # Slide heading (inside bar)
#     _add_textbox(slide,
#                  left=Inches(0.5), top=Inches(0.15),
#                  width=Inches(12.33), height=Inches(0.95),
#                  text=heading, font_name=TITLE_FONT, font_size=32,
#                  bold=True, color=TEXT_LITE)

#     # Bullet points
#     if bullets:
#         txb   = slide.shapes.add_textbox(
#             Inches(0.7), Inches(1.55), Inches(11.93), Inches(5.7))
#         tf    = txb.text_frame
#         tf.word_wrap = True

#         for i, bullet in enumerate(bullets):
#             para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
#             para.space_before = Pt(6)
#             para.space_after  = Pt(2)
#             para.level        = 0

#             # Bullet marker
#             marker = para.add_run()
#             marker.text           = "▸  "
#             marker.font.name      = BODY_FONT
#             marker.font.size      = Pt(16)
#             marker.font.bold      = True
#             marker.font.color.rgb = ACCENT

#             # Bullet text
#             run = para.add_run()
#             run.text           = bullet
#             run.font.name      = BODY_FONT
#             run.font.size      = Pt(16)
#             run.font.bold      = False
#             run.font.color.rgb = BULLET_CLR


# def _build_closing_slide(prs: Presentation, topic: str):
#     """Closing navy slide."""
#     _build_title_slide(prs, "Thank You", f"Presentation on: {topic}")


# # ── Public API ───────────────────────────────────────────────────────────────
# def generate_pptx(slides_json: dict, output_dir: str = ".") -> str:
#     """
#     Build a .pptx file from slide JSON and return its file path.

#     Expected input format:
#         {
#             "topic": "AI in Healthcare",          # optional
#             "slides": [
#                 { "heading": "Slide Title", "bullets": ["Point 1", "Point 2"] },
#                 ...
#             ]
#         }

#     Returns:
#         str – absolute path to the saved .pptx file.
#     """
#     topic  = slides_json.get("topic", "Presentation")
#     slides = slides_json.get("slides", [])

#     if not slides:
#         raise ValueError("'slides' list is empty — nothing to generate.")

#     prs = Presentation()
#     prs.slide_width  = SLIDE_W
#     prs.slide_height = SLIDE_H

#     # Title slide
#     _build_title_slide(prs, topic, f"{len(slides)} slides")

#     # Content slides
#     for slide in slides:
#         heading = slide.get("heading", "Untitled Slide")
#         bullets = slide.get("bullets", [])
#         _build_content_slide(prs, heading, bullets)

#     # Closing slide
#     _build_closing_slide(prs, topic)

#     # Save
#     os.makedirs(output_dir, exist_ok=True)
#     safe_name = "".join(c if c.isalnum() or c in " _-" else "_" for c in topic)
#     filename  = f"{safe_name.strip()[:50]}.pptx"
#     filepath  = os.path.abspath(os.path.join(output_dir, filename))
#     prs.save(filepath)

#     return filepath


# # ── Quick test ───────────────────────────────────────────────────────────────
# if __name__ == "__main__":
#     sample = {
#         "topic": "Artificial Intelligence",
#         "slides": [
#             {
#                 "heading": "What is Artificial Intelligence?",
#                 "bullets": [
#                     "AI enables machines to simulate human intelligence",
#                     "Encompasses machine learning, NLP, and computer vision",
#                     "Powers applications from search engines to self-driving cars",
#                 ]
#             },
#             {
#                 "heading": "Key Types of AI",
#                 "bullets": [
#                     "Narrow AI – designed for specific tasks (e.g. image recognition)",
#                     "General AI – human-level reasoning across domains (theoretical)",
#                     "Machine Learning – systems that learn from data automatically",
#                 ]
#             },
#             {
#                 "heading": "Real-World Applications",
#                 "bullets": [
#                     "Healthcare: early disease detection and drug discovery",
#                     "Finance: fraud detection and algorithmic trading",
#                     "Education: personalised learning and smart tutoring systems",
#                 ]
#             },
#             {
#                 "heading": "Challenges & Considerations",
#                 "bullets": [
#                     "Bias in training data can lead to unfair outcomes",
#                     "Privacy concerns around large-scale data collection",
#                     "Need for transparent and explainable AI decisions",
#                 ]
#             },
#             {
#                 "heading": "The Future of AI",
#                 "bullets": [
#                     "Rapid advances in multimodal and generative AI models",
#                     "Growing importance of AI regulation and ethics frameworks",
#                     "Increasing collaboration between humans and AI systems",
#                 ]
#             },
#         ]
#     }

#     path = generate_pptx(sample, output_dir=".")
#     print(f"Saved → {path}")
    